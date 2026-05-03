"""Generate StockPulse.pptx"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette ──────────────────────────────────────────────────────────
BG      = RGBColor(0x0a, 0x0e, 0x1a)
SURFACE = RGBColor(0x11, 0x18, 0x27)
CARD    = RGBColor(0x1a, 0x22, 0x35)
BORDER  = RGBColor(0x1e, 0x2d, 0x45)
BLUE    = RGBColor(0x3b, 0x82, 0xf6)
PURPLE  = RGBColor(0x8b, 0x5c, 0xf6)
CYAN    = RGBColor(0x06, 0xb6, 0xd4)
GREEN   = RGBColor(0x10, 0xb9, 0x81)
YELLOW  = RGBColor(0xf5, 0x9e, 0x0b)
RED     = RGBColor(0xef, 0x44, 0x44)
WHITE   = RGBColor(0xff, 0xff, 0xff)
MUTED   = RGBColor(0x64, 0x74, 0x8b)
LIGHT   = RGBColor(0xe2, 0xe8, 0xf0)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb=None, alpha=None, line_rgb=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_w
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=LIGHT, align=PP_ALIGN.LEFT,
             wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def slide_bg(slide, rgb=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb

def accent_bar(slide, color=BLUE, height=Pt(4)):
    add_rect(slide, 0, 0, W, height, fill_rgb=color)

def section_label(slide, text, y=Inches(0.55)):
    add_text(slide, text.upper(), Inches(0.7), y, Inches(6), Inches(0.35),
             size=10, bold=True, color=BLUE)

def slide_title(slide, text, y=Inches(0.9)):
    add_text(slide, text, Inches(0.7), y, Inches(11.9), Inches(1.0),
             size=36, bold=True, color=WHITE)

def bullet_box(slide, lines, x, y, w, h, icon_color=BLUE):
    """Render a simple card with bullet lines."""
    add_rect(slide, x, y, w, h, fill_rgb=CARD, line_rgb=BORDER, line_w=Pt(1))
    txBox = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.18),
                                      w - Inches(0.36), h - Inches(0.28))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(13)
        run.font.color.rgb = LIGHT
        run.font.name = "Calibri"

def stat_card(slide, number, label, x, y, w=Inches(2.8), h=Inches(1.4), num_color=BLUE):
    add_rect(slide, x, y, w, h, fill_rgb=CARD, line_rgb=BORDER, line_w=Pt(1))
    # top accent stripe
    add_rect(slide, x, y, w, Pt(3), fill_rgb=num_color)
    add_text(slide, number, x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), Inches(0.65),
             size=32, bold=True, color=num_color, align=PP_ALIGN.LEFT)
    add_text(slide, label, x + Inches(0.15), y + Inches(0.75), w - Inches(0.3), Inches(0.55),
             size=11, color=MUTED, align=PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)

# gradient-ish top glow via large soft rect
add_rect(sl, Inches(2), Inches(0), Inches(9.33), Inches(4.5), fill_rgb=RGBColor(0x0d, 0x1a, 0x38))

# main title
add_text(sl, "📈  StockPulse",
         Inches(0.5), Inches(1.6), Inches(12.33), Inches(1.6),
         size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "股市新聞爬蟲 · 情緒分析 · Dashboard · LINE Bot",
         Inches(0.5), Inches(3.3), Inches(12.33), Inches(0.7),
         size=20, color=MUTED, align=PP_ALIGN.CENTER)

# tags row
tags = [("Python", BLUE), ("PostgreSQL", PURPLE), ("Flask", CYAN),
        ("NVIDIA NIM", GREEN), ("LINE Bot", RGBColor(0x00, 0xb9, 0x00))]
tx = Inches(1.6)
for label, col in tags:
    add_rect(sl, tx, Inches(4.25), Inches(1.75), Inches(0.45),
             fill_rgb=RGBColor(0x0d, 0x12, 0x22),
             line_rgb=col, line_w=Pt(1))
    add_text(sl, label, tx + Inches(0.06), Inches(4.26), Inches(1.63), Inches(0.42),
             size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    tx += Inches(1.95)

add_text(sl, "2026  ·  Brian",
         Inches(0.5), Inches(6.9), Inches(12.33), Inches(0.4),
         size=12, color=MUTED, align=PP_ALIGN.CENTER)

accent_bar(sl, BLUE)

# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — 系統概覽
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
accent_bar(sl)
section_label(sl, "Overview")
slide_title(sl, "系統概覽")

# arch boxes
boxes = [
    ("🌐\n新聞來源", "Yahoo / 鉅亨 / NStock", BLUE),
    ("🕷️\n分級爬蟲", "3 來源 × 4 級別", PURPLE),
    ("🧠\nAI 情緒分析", "NVIDIA NIM  1–9 分", CYAN),
    ("🗄️\nPostgreSQL", "統一資料庫", GREEN),
    ("📊\n展示層", "Dashboard + LINE", YELLOW),
]
box_w = Inches(2.2)
box_h = Inches(1.55)
start_x = Inches(0.4)
gap = Inches(0.3)
by = Inches(2.1)

for i, (icon_lbl, sub, col) in enumerate(boxes):
    bx = start_x + i * (box_w + gap)
    add_rect(sl, bx, by, box_w, box_h, fill_rgb=CARD, line_rgb=col, line_w=Pt(1.5))
    add_rect(sl, bx, by, box_w, Pt(3), fill_rgb=col)
    add_text(sl, icon_lbl, bx + Inches(0.1), by + Inches(0.1),
             box_w - Inches(0.2), Inches(0.75), size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, sub, bx + Inches(0.1), by + Inches(0.9),
             box_w - Inches(0.2), Inches(0.55), size=11, color=MUTED, align=PP_ALIGN.CENTER)
    if i < len(boxes) - 1:
        ax = bx + box_w + Inches(0.04)
        add_text(sl, "→", ax, by + Inches(0.5), gap - Inches(0.04), Inches(0.5),
                 size=18, color=BLUE, align=PP_ALIGN.CENTER)

# stat cards row
stat_data = [("3", "新聞來源", BLUE), ("1,949", "覆蓋股票數", PURPLE), ("2,978+", "資料庫新聞筆數", GREEN)]
sx = Inches(0.85)
for num, lbl, col in stat_data:
    stat_card(sl, num, lbl, sx, Inches(4.1), w=Inches(3.8), h=Inches(1.55), num_color=col)
    sx += Inches(4.05)

# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — 爬蟲架構
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
accent_bar(sl)
section_label(sl, "Crawlers")
slide_title(sl, "爬蟲架構")

# Left col: 3 phases
phases = [
    (BLUE,   "① Phase 1 — 抓新聞列表",
             "curl_cffi 模擬真實瀏覽器 TLS 指紋，繞過反爬蟲\n非同步並發 · 指數退避重試 · 斷點續爬"),
    (PURPLE, "② Phase 2 — 抓新聞內文",
             "fetch_content.py 補抓完整文章內容\n失敗 URL 記錄，可重跑"),
    (GREEN,  "③ Phase 3 — AI 情緒分析",
             "NVIDIA NIM API（Llama 系列）\n給每篇新聞打 1–9 分情緒評分"),
]
py = Inches(2.0)
for col, title, desc in phases:
    add_rect(sl, Inches(0.5), py, Inches(6.0), Inches(1.3), fill_rgb=CARD, line_rgb=col, line_w=Pt(1.5))
    add_rect(sl, Inches(0.5), py, Pt(4), Inches(1.3), fill_rgb=col)
    add_text(sl, title, Inches(0.75), py + Inches(0.1), Inches(5.6), Inches(0.45),
             size=13, bold=True, color=WHITE)
    add_text(sl, desc, Inches(0.75), py + Inches(0.52), Inches(5.6), Inches(0.68),
             size=11, color=MUTED)
    py += Inches(1.45)

# Right col: tiers
tiers = [
    (BLUE,   "🔥 熱門股", "974 支 · 每日抓 3 天內新聞"),
    (CYAN,   "📊 中間股", "487 支 · 每日抓 3 天內新聞"),
    (YELLOW, "📉 偏下股", "293 支 · 每日抓 7 天內新聞"),
    (PURPLE, "💤 稀少股", "195 支 · 每日抓 7 天內新聞"),
]
add_text(sl, "股票分級策略",
         Inches(7.0), Inches(1.9), Inches(5.8), Inches(0.4),
         size=13, bold=True, color=MUTED)
ty = Inches(2.35)
for col, title, desc in tiers:
    add_rect(sl, Inches(7.0), ty, Inches(5.8), Inches(1.05), fill_rgb=CARD, line_rgb=col, line_w=Pt(1.5))
    add_text(sl, title, Inches(7.2), ty + Inches(0.1), Inches(5.4), Inches(0.4),
             size=13, bold=True, color=col)
    add_text(sl, desc, Inches(7.2), ty + Inches(0.5), Inches(5.4), Inches(0.45),
             size=11, color=MUTED)
    ty += Inches(1.15)

# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — 三大新聞來源
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
accent_bar(sl)
section_label(sl, "Data Sources")
slide_title(sl, "三大新聞來源")

sources = [
    ("🟣  Yahoo Finance",  PURPLE, ["crawler_hot / mid / lower / rare", "fetch_content.py", "analyze_sentiment.py", "→ 6 支腳本"]),
    ("🔵  鉅亨網 Cnyes",   BLUE,   ["crawler_hot / mid / lower / rare", "fetch_content.py", "analyze_sentiment.py", "→ 6 支腳本"]),
    ("🟢  NStock",         GREEN,  ["crawler_hot / mid / lower / rare", "fetch_content.py", "analyze_sentiment.py", "→ 6 支腳本 + 工具"]),
]
sx = Inches(0.45)
for title, col, lines in sources:
    add_rect(sl, sx, Inches(2.05), Inches(4.1), Inches(3.2), fill_rgb=CARD, line_rgb=col, line_w=Pt(1.5))
    add_rect(sl, sx, Inches(2.05), Inches(4.1), Pt(3), fill_rgb=col)
    add_text(sl, title, sx + Inches(0.15), Inches(2.15), Inches(3.8), Inches(0.6),
             size=16, bold=True, color=col)
    ly = Inches(2.8)
    for line in lines:
        add_text(sl, line, sx + Inches(0.15), ly, Inches(3.8), Inches(0.5),
                 size=12, color=LIGHT if "→" not in line else col, bold="→" in line)
        ly += Inches(0.52)
    sx += Inches(4.35)

# bottom banner
add_rect(sl, Inches(0.45), Inches(5.6), Inches(12.43), Inches(1.1), fill_rgb=CARD, line_rgb=BLUE, line_w=Pt(1))
add_text(sl, "⏰  自動化排程 — run_crawlers_daily.bat",
         Inches(0.65), Inches(5.68), Inches(7), Inches(0.45), size=14, bold=True, color=WHITE)
add_text(sl, "開機等待網路 → 清除舊進度 → 3 來源 × 3 Phase → 記錄日誌",
         Inches(0.65), Inches(6.12), Inches(8), Inches(0.4), size=12, color=MUTED)

# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — 資料庫結構
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
accent_bar(sl)
section_label(sl, "Database")
slide_title(sl, "資料庫結構（PostgreSQL）")

# Table header
add_rect(sl, Inches(0.5), Inches(2.05), Inches(7.5), Inches(0.45),
         fill_rgb=RGBColor(0x1e, 0x3a, 0x5f), line_rgb=BLUE, line_w=Pt(1))
add_text(sl, "🗄️  yahoo / cnyes / nstock_stock_news",
         Inches(0.65), Inches(2.08), Inches(7.2), Inches(0.38), size=12, bold=True, color=CYAN)

rows_main = [
    ("stock_id",       "VARCHAR",   "股票代碼"),
    ("title",          "TEXT",      "新聞標題"),
    ("link",           "TEXT",      "原始網址"),
    ("content",        "TEXT",      "新聞內文"),
    ("sentiment_score","INTEGER",   "情緒分數 1–9"),
    ("fetched_at",     "TIMESTAMP", "抓取時間"),
    ("fetched_date",   "DATE",      "抓取日期（索引）"),
]
ry = Inches(2.5)
for col_name, col_type, col_desc in rows_main:
    add_rect(sl, Inches(0.5), ry, Inches(7.5), Inches(0.47),
             fill_rgb=CARD, line_rgb=BORDER, line_w=Pt(0.5))
    add_text(sl, col_name, Inches(0.65), ry + Inches(0.08), Inches(2), Inches(0.35),
             size=11, color=CYAN, font_name="Courier New")
    add_text(sl, col_type, Inches(2.8), ry + Inches(0.08), Inches(1.5), Inches(0.35),
             size=11, color=PURPLE)
    add_text(sl, col_desc, Inches(4.4), ry + Inches(0.08), Inches(3.4), Inches(0.35),
             size=11, color=MUTED)
    ry += Inches(0.47)

# stock_mapping
add_rect(sl, Inches(8.3), Inches(2.05), Inches(4.6), Inches(0.45),
         fill_rgb=RGBColor(0x1e, 0x3a, 0x5f), line_rgb=GREEN, line_w=Pt(1))
add_text(sl, "📋  stock_mapping",
         Inches(8.45), Inches(2.08), Inches(4.3), Inches(0.38), size=12, bold=True, color=GREEN)

for col_name, col_type, col_desc in [("stock_id", "VARCHAR PK", "股票代碼"),
                                      ("stock_name", "VARCHAR", "股票中文名稱")]:
    add_rect(sl, Inches(8.3), ry - Inches(2.82), Inches(4.6), Inches(0.47),
             fill_rgb=CARD, line_rgb=BORDER, line_w=Pt(0.5))
    add_text(sl, col_name, Inches(8.45), ry - Inches(2.74), Inches(1.6), Inches(0.35),
             size=11, color=CYAN, font_name="Courier New")
    add_text(sl, col_type, Inches(10.15), ry - Inches(2.74), Inches(1.3), Inches(0.35),
             size=11, color=PURPLE)
    ry += Inches(0.47)

# note card
add_rect(sl, Inches(8.3), Inches(3.15), Inches(4.6), Inches(1.4),
         fill_rgb=CARD, line_rgb=GREEN, line_w=Pt(1))
add_text(sl, "🔗  跨源合併查詢",
         Inches(8.45), Inches(3.22), Inches(4.3), Inches(0.45), size=13, bold=True, color=GREEN)
add_text(sl, "UNION ALL 合併所有來源\nDISTINCT ON 去重\n避免同篇新聞重複顯示",
         Inches(8.45), Inches(3.65), Inches(4.3), Inches(0.75), size=11, color=MUTED)


# ════════════════════════════════════════════════════════════════════
# SLIDE 6 — Dashboard
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
accent_bar(sl)
section_label(sl, "Frontend")
slide_title(sl, "Dashboard 功能")

kpi = [
    ("2,472", "本週總新聞", PURPLE),
    ("6.09",  "平均情緒",   RGBColor(0xec, 0x48, 0x99)),
    ("913",   "正面新聞",   BLUE),
    ("130",   "負面新聞",   YELLOW),
]
kx = Inches(0.35)
for num, lbl, col in kpi:
    stat_card(sl, num, lbl, kx, Inches(2.05), w=Inches(3.05), h=Inches(1.45), num_color=col)
    kx += Inches(3.23)

feature_cards = [
    (Inches(0.35), "🔥  TOP 5 熱門股",
     "1. 華邦電 (23 則)\n2. 旺宏 (22 則)\n3. 力積電 (21 則)\n4. 台積電 (21 則)\n5. 聯電 (21 則)", BLUE),
    (Inches(4.6),  "📊  情緒趨勢折線圖",
     "過去 7 天每日平均情緒\nChart.js v4.4 · 橙色主題\n趨勢：5.58 → 6.60 📈", CYAN),
    (Inches(8.85), "⚡  即時更新",
     "每 30 秒自動刷新\n最新 20 則新聞串流\n近 30 天趨勢圖", GREEN),
]
for fx, ft, fd, fc in feature_cards:
    add_rect(sl, fx, Inches(3.75), Inches(4.1), Inches(2.8), fill_rgb=CARD, line_rgb=fc, line_w=Pt(1.5))
    add_rect(sl, fx, Inches(3.75), Inches(4.1), Pt(3), fill_rgb=fc)
    add_text(sl, ft, fx + Inches(0.15), Inches(3.85), Inches(3.8), Inches(0.55),
             size=14, bold=True, color=fc)
    add_text(sl, fd, fx + Inches(0.15), Inches(4.45), Inches(3.8), Inches(2.0),
             size=12, color=LIGHT)

# ════════════════════════════════════════════════════════════════════
# SLIDE 7 — LINE Bot
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
accent_bar(sl, RGBColor(0x00, 0xc3, 0x00))
section_label(sl, "LINE Bot")
slide_title(sl, "LINE Bot 指令")

cmds = [
    (BLUE,   "🔍  查詢 [股票名稱]", "查詢特定股票最新新聞"),
    (YELLOW, "🔥  熱門",            "最活躍 10 檔股票排行"),
    (CYAN,   "🆕  最新",            "最新 5 則新聞"),
    (GREEN,  "📈  正面 / 負面",      "依情緒分數過濾（7-9 / 1-3）"),
    (PURPLE, "📊  一周 / 本周",      "最近 7 天新聞分析摘要"),
]
cy = Inches(2.05)
for col, title, desc in cmds:
    add_rect(sl, Inches(0.5), cy, Inches(6.0), Inches(0.95), fill_rgb=CARD, line_rgb=col, line_w=Pt(1.5))
    add_rect(sl, Inches(0.5), cy, Pt(4), Inches(0.95), fill_rgb=col)
    add_text(sl, title, Inches(0.75), cy + Inches(0.08), Inches(5.6), Inches(0.42),
             size=13, bold=True, color=col)
    add_text(sl, desc, Inches(0.75), cy + Inches(0.5), Inches(5.6), Inches(0.38),
             size=11, color=MUTED)
    cy += Inches(1.05)

# Sample reply card
add_rect(sl, Inches(7.1), Inches(2.05), Inches(5.8), Inches(4.15),
         fill_rgb=RGBColor(0x06, 0x12, 0x22), line_rgb=GREEN, line_w=Pt(1.5))
add_text(sl, "📱  本周回應範例",
         Inches(7.25), Inches(2.12), Inches(5.5), Inches(0.45), size=13, bold=True, color=GREEN)
sample = (
    "📊 最近一周新聞分析\n\n"
    "📰 總新聞數：2472 則\n"
    "📈 平均情緒：6.09\n"
    "📈 正面新聞：913 則\n"
    "📉 負面新聞：130 則\n\n"
    "🔥 本周最熱門股票：\n"
    "1. 華邦電 (23 則)\n"
    "2. 旺宏 (22 則)\n"
    "3. 台積電 (21 則)"
)
add_text(sl, sample, Inches(7.25), Inches(2.6), Inches(5.5), Inches(3.5),
         size=12, color=LIGHT, font_name="Courier New")

# ════════════════════════════════════════════════════════════════════
# SLIDE 8 — 技術棧
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
accent_bar(sl)
section_label(sl, "Tech Stack")
slide_title(sl, "技術棧")

crawler_tech = [
    ("🐍", "Python 3.14"),
    ("🌐", "curl_cffi"),
    ("🍲", "BeautifulSoup4"),
    ("⚡", "asyncio / asyncpg"),
    ("🧠", "NVIDIA NIM API"),
    ("🤖", "OpenAI SDK"),
    ("🌶️", "Flask"),
    ("🐘", "PostgreSQL"),
]
frontend_tech = [
    ("📊", "Chart.js v4.4"),
    ("💻", "Vanilla JS + CSS3"),
    ("💬", "LINE Bot SDK"),
    ("🔗", "ngrok (Webhook)"),
    ("🔐", "python-dotenv"),
    ("📋", "Windows Task Scheduler"),
]

add_text(sl, "爬蟲 / 後端", Inches(0.5), Inches(2.0), Inches(6.2), Inches(0.4),
         size=12, bold=True, color=MUTED)
tx, ty = Inches(0.5), Inches(2.45)
for icon, name in crawler_tech:
    add_rect(sl, tx, ty, Inches(2.9), Inches(0.55), fill_rgb=CARD, line_rgb=BORDER, line_w=Pt(1))
    add_text(sl, f"{icon}  {name}", tx + Inches(0.12), ty + Inches(0.1),
             Inches(2.65), Inches(0.38), size=12, color=LIGHT)
    tx += Inches(3.1)
    if tx > Inches(4):
        tx = Inches(0.5)
        ty += Inches(0.65)

add_text(sl, "前端 / 整合", Inches(6.85), Inches(2.0), Inches(6.0), Inches(0.4),
         size=12, bold=True, color=MUTED)
tx, ty = Inches(6.85), Inches(2.45)
for icon, name in frontend_tech:
    add_rect(sl, tx, ty, Inches(2.9), Inches(0.55), fill_rgb=CARD, line_rgb=BORDER, line_w=Pt(1))
    add_text(sl, f"{icon}  {name}", tx + Inches(0.12), ty + Inches(0.1),
             Inches(2.65), Inches(0.38), size=12, color=LIGHT)
    tx += Inches(3.1)
    if tx > Inches(10.9):
        tx = Inches(6.85)
        ty += Inches(0.65)

# Security note
add_rect(sl, Inches(0.5), Inches(5.85), Inches(12.33), Inches(1.0),
         fill_rgb=CARD, line_rgb=PURPLE, line_w=Pt(1))
add_text(sl, "🔒  安全性",
         Inches(0.65), Inches(5.9), Inches(2), Inches(0.4), size=13, bold=True, color=PURPLE)
add_text(sl, "API Keys 改用 .env 環境變數  ·  git filter-repo 清除歷史洩漏  ·  .gitignore 防再次 commit",
         Inches(2.6), Inches(5.92), Inches(10.0), Inches(0.75), size=12, color=MUTED)

# ════════════════════════════════════════════════════════════════════
# SLIDE 9 — Roadmap & Thank You
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
accent_bar(sl, PURPLE)
section_label(sl, "Roadmap")
slide_title(sl, "未來規劃")

short_term = [
    "週對比功能（本周 vs 上周）",
    "行業分類分析",
    "移動端 Dashboard 優化",
    "LINE Bot 圖片回覆",
]
long_term = [
    "自訂時間範圍查詢（3/14/30 天）",
    "情緒異常波動預警通知",
    "匯出 PDF 報告",
    "郵件訂閱功能",
]

for col_x, col, title, items in [
    (Inches(0.5),  BLUE,   "🔥  短期優化",  short_term),
    (Inches(6.85), PURPLE, "🚀  長期擴展",  long_term),
]:
    add_rect(sl, col_x, Inches(2.05), Inches(6.0), Inches(3.5),
             fill_rgb=CARD, line_rgb=col, line_w=Pt(1.5))
    add_rect(sl, col_x, Inches(2.05), Inches(6.0), Pt(3), fill_rgb=col)
    add_text(sl, title, col_x + Inches(0.15), Inches(2.12), Inches(5.7), Inches(0.5),
             size=15, bold=True, color=col)
    iy = Inches(2.72)
    for item in items:
        add_text(sl, f"→  {item}", col_x + Inches(0.15), iy, Inches(5.7), Inches(0.5),
                 size=12, color=LIGHT)
        iy += Inches(0.58)

# Thank you banner
add_rect(sl, Inches(0.5), Inches(5.9), Inches(12.33), Inches(1.2),
         fill_rgb=RGBColor(0x0d, 0x1a, 0x38), line_rgb=BLUE, line_w=Pt(1))
add_text(sl, "🎉  謝謝觀看   ·   StockPulse  ·  2026",
         Inches(0.5), Inches(6.0), Inches(12.33), Inches(0.9),
         size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Save ────────────────────────────────────────────────────────────────────
out = r"d:\StockPulse\StockPulse簡報.pptx"
prs.save(out)
print(f"OK Saved: {out}")
